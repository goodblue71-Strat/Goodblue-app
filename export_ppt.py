"""
PowerPoint export skeleton for Strategy Copilot (Step 1)

- Requires: python-pptx (pip install python-pptx)
- Builds a polished 16:9 deck from your session_state
- Slides: Title, Agenda, Executive Snapshot, SWOT, Ansoff 2x2, Benchmark table, Top-5 Recommendations (Impact×Effort grid), Appendix

Usage in Streamlit (Export step):

    from export_ppt import build_ppt_from_state
    bio, fname = build_ppt_from_state(state)
    st.download_button("Download PPTX", data=bio.getvalue(), file_name=fname, mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

This module is defensive: missing sections are skipped gracefully.
"""
from __future__ import annotations
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from typing import Any, Dict, List, Optional
from pptx.dml.color import RGBColor
from io import BytesIO
from datetime import datetime
import numpy as np
import streamlit as st
import pandas as pd

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.8)
TITLE_SIZE = Pt(36)
SUBTITLE_SIZE = Pt(18)
H2_SIZE = Pt(24)
BODY_SIZE = Pt(14)
MONO_SIZE = Pt(10)

COLOR_PRIMARY = RGBColor(30, 64, 175)    # blue-700
COLOR_ACCENT  = RGBColor(16, 185, 129)   # emerald-500
COLOR_DARK    = RGBColor(17, 24, 39)     # gray-900
COLOR_MED     = RGBColor(75, 85, 99)     # gray-600
COLOR_LIGHT   = RGBColor(229, 231, 235)  # gray-200

def _members(names):
    # only include names that exist in this python-pptx version
    return {getattr(PP_PLACEHOLDER, n) for n in names if hasattr(PP_PLACEHOLDER, n)}

TITLE_TYPES   = _members(["TITLE", "CENTER_TITLE", "VERTICAL_TITLE"])
BODYISH_TYPES = _members(["BODY", "VERTICAL_BODY", "OBJECT", "TABLE", "CHART", "PICTURE"])
META_TYPES    = _members(["DATE", "SLIDE_NUMBER", "FOOTER", "HEADER"])

BODYISH = {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}
if hasattr(PP_PLACEHOLDER, "VERTICAL_BODY"):
    BODYISH.add(PP_PLACEHOLDER.VERTICAL_BODY)

def has_body(layout) -> bool:
        return any(getattr(ph, "placeholder_format", None) and ph.placeholder_format.type in BODYISH_TYPES
            for ph in layout.placeholders)

def has_title(layout) -> bool:
    return any(getattr(ph, "placeholder_format", None)
               and ph.placeholder_format.type in TITLE_TYPES
               for ph in layout.placeholders)



def is_blank(layout) -> bool:
    types = {ph.placeholder_format.type for ph in layout.placeholders}
    return len(types) == 0 or types.issubset(META_TYPES)

def _add_title(prs, title, subtitle=None):
    st.write(title)
    # helper: does a layout truly have a title placeholder?
    def has_title(layout):
        for ph in layout.placeholders:
            pf = getattr(ph, "placeholder_format", None)
            if pf and pf.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                return True
        return False

    # choose a layout with a real title placeholder (fallback to first)
    layout = next((l for l in prs.slide_layouts if has_title(l)), prs.slide_layouts[0])
    slide = prs.slides.add_slide(layout)

    # set title (guard if None)
    t = slide.shapes.title
    if t is not None:
        t.text = title
        title_top = t.top
        title_height = t.height
    else:
        margin = Inches(0.6)
        title_top = Inches(0.6)
        title_height = Inches(1.0)
        tb = slide.shapes.add_textbox(margin, title_top, prs.slide_width - 2*margin, title_height)
        tb.text_frame.clear()
        tb.text_frame.paragraphs[0].add_run().text = title

    # try to use a real SUBTITLE placeholder
    placed_subtitle = False
    if subtitle:
        for ph in slide.placeholders:
            pf = getattr(ph, "placeholder_format", None)
            if pf and pf.type == PP_PLACEHOLDER.SUBTITLE:
                ph.text = subtitle
                placed_subtitle = True
                break

        if not placed_subtitle:
            margin = Inches(0.6)
            top = title_top + title_height + Inches(0.2)
            sub_box = slide.shapes.add_textbox(margin, top, prs.slide_width - 2*margin, Inches(0.9))
            stf = sub_box.text_frame
            stf.clear()
            p2 = stf.paragraphs[0]
            r2 = p2.add_run()
            r2.text = subtitle

    # Accent bar (optional)
    margin = Inches(0.6)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, margin, Inches(0.6), Inches(2.2), Inches(0.15))
    # rect.fill.solid(); rect.fill.fore_color.rgb = COLOR_PRIMARY
    return slide

def _add_heading(slide, text: str):
    title_shape = slide.shapes.title
    if title_shape is not None:
        tf = title_shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = text
        return(title_shape)
    # Fallback: make a heading textbox
    box = slide.shapes.add_textbox(MARGIN, MARGIN, W - 2*MARGIN, Inches(0.6))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = H2_SIZE
    run.font.bold = True
    run.font.color.rgb = COLOR_DARK
    return box


def _add_bullets(slide, left, top, width, height, items: List[str]):
    body = next((ph for ph in slide.placeholders
                 if getattr(ph, "placeholder_format", None)
                 and ph.placeholder_format.type in BODYISH), None)

    # Get a text_frame from body if present; otherwise create a textbox
    if body is not None:
        used_shape = body
        tf = body.text_frame
        force_bullets = False
    else:
        used_shape = slide.shapes.add_textbox(left, top, width, height)
        tf = used_shape.text_frame
        force_bullets = True  # textboxes often don’t bullet by default

    tf.word_wrap = True
    tf.clear()

    for i, item in enumerate(items or []):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if force_bullets:
            try:
                p._element.get_or_add_pPr().get_or_add_buChar()  # force symbol bullets
            except Exception:
                pass
        p.level = 0
        p.space_after = Pt(6)

        r = p.add_run()
        r.text = str(item)
        r.font.size = Pt(BODY_SIZE.pt + 3) if hasattr(BODY_SIZE, "pt") else Pt(BODY_SIZE + 3)
        r.font.color.rgb = COLOR_DARK
    return used_shape


def _grid(slide, left, top, width, height):
    # Draw outer rect and cross-lines; return quadrant rects
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = COLOR_MED
    # Vertical and horizontal lines
    slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, left + width/2, top, Inches(0.0), height).line.color.rgb = COLOR_LIGHT
    slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, left, top + height/2, width, Inches(0.0)).line.color.rgb = COLOR_LIGHT
    # Return quadrants as (l, t, w, h)
    return [
        (left, top, width/2, height/2),                # Q1 TL
        (left + width/2, top, width/2, height/2),      # Q2 TR
        (left, top + height/2, width/2, height/2),     # Q3 BL
        (left + width/2, top + height/2, width/2, height/2),  # Q4 BR
    ]


def _add_small_label(slide, text: str, left, top,angle_deg: float = 0):
    b = slide.shapes.add_textbox(left, top, Inches(2), Inches(0.3))
    b.rotation = angle_deg % 360
    tf = b.text_frame; tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # keep text vertically centered
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text; r.font.size = Pt(11); r.font.color.rgb = COLOR_MED

# ---------------------------- Slide builders ----------------------------

def slide_agenda(prs: Presentation, items: Optional[List[str]] = None):
    #blank = next((l for l in prs.slide_layouts if len(l.placeholders) == 0), prs.slide_layouts[0])
    #slide = prs.slides.add_slide(blank)
    layout = next((l for l in prs.slide_layouts if has_body(l)), prs.slide_layouts[0])
    slide = prs.slides.add_slide(layout)
    _add_heading(slide, "Agenda")
    _add_bullets(slide, MARGIN, Inches(2.0), W - 2*MARGIN, Inches(5.0), items or [
        "Inputs & Goals",
        "Framework Insights",
        "Recommendations",
        "Next Steps",
    ])
    return slide

def slide_exec_snapshot(prs: Presentation, bullets: List[str]):
    layout = next((l for l in prs.slide_layouts if has_body(l)), prs.slide_layouts[0])
    slide = prs.slides.add_slide(layout)
    _add_heading(slide, "Executive Snapshot")
    _add_bullets(slide, MARGIN, Inches(1.2), W - 2*MARGIN, Inches(5.0), bullets)
    return slide

def _top_factor(cs: Dict[str, Any], keys: List[str]) -> str:
    """Return the first factor for the first matching key (handles key variants)."""
    for k in keys:
        v = cs.get(k)
        if isinstance(v, list) and v:
            return _first(v)
    return ""

def ind_to_long_records(ind_raw):
    """Accepts a list or {'industries': [...]} and returns long records."""
    if isinstance(ind_raw, list):
        industries = ind_raw
    elif isinstance(ind_raw, dict):
        industries = ind_raw.get("industries", [])
    else:
        industries = []

    records = []
    for item in industries:
        industry = item.get("industry_vertical_name", "")
        tam = item.get("TAM", "")
        cs = item.get("Critical_success_category", {}) or {}

        for category, factors in cs.items():
            if isinstance(factors, list):
                for rank, factor in enumerate(factors, start=1):
                    records.append({
                        "Industry": industry,
                        "TAM (B$)": tam,
                        "Category": category,
                        "Rank": rank,
                        "Factor": str(factor),
                    })
            else:
                # tolerate non-list shapes
                records.append({
                    "Industry": industry,
                    "TAM (B$)": tam,
                    "Category": category,
                    "Rank": 1,
                    "Factor": str(factors),
                })
    return records

"""
def slide_ind(prs, ind_raw):
    records = ind_to_long_records(ind_raw)
    lines=[]
    for r in records:
        line = "\t".join(str(r))
        lines.append(line)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _add_heading(slide, "INDUSTRY ANALYSIS")
    lines_per_slide=20
    for start in range(0, len(lines), lines_per_slide):
        chunk = "\n".join(lines[start:start+lines_per_slide])
        box = slide.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9.0), Inches(5.0))
        t = box.text_frame
        t.clear()
        t.text = chunk
        for p in t.paragraphs:
            p.font.name = "Courier New"   # columns line up better
            p.font.size = Pt(10)
    return slide
"""

def slide_ind(prs, ind_raw):
    records = ind_to_long_records(ind_raw)
    if not records:
        return

    # Get headers from first record
    headers = list(records[0].keys())
    rows, cols = len(records) + 1, len(headers)

    # Add a blank slide
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    shapes = slide.shapes

    # Add title
    _add_heading(slide, "INDUSTRY ANALYSIS")

    # Add table
    table = shapes.add_table(rows, cols, Inches(0.5), Inches(1.0), Inches(9.0), Inches(5.0)).table

    # Format headers
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(12)

    # Fill rows
    for row_idx, record in enumerate(records, start=1):
        for col_idx, header in enumerate(headers):
            table.cell(row_idx, col_idx).text = str(record.get(header, ""))

    return slide
 
def slide_swot(prs: Presentation, swot: Dict[str, List[str]]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _add_heading(slide, "SWOT")
    box_w = (W - 3*MARGIN) / 2
    box_h = (H - 2*MARGIN - Inches(1.0)) / 2
    x1, x2 = 0.5*MARGIN, 0.5*MARGIN + box_w + MARGIN
    y1, y2 = Inches(1.5), Inches(1.5) + box_h + 0.5*MARGIN

    def cell(title, items, x, y):
        title_box = slide.shapes.add_textbox(x, y, box_w, Inches(0.35))
        tf = title_box.text_frame; tf.clear(); p = tf.paragraphs[0]; r = p.add_run(); r.text = title; r.font.bold = True; r.font.size = Pt(16); r.font.color.rgb = COLOR_PRIMARY
        _add_bullets(slide, x, y + Inches(0.4), box_w, box_h - Inches(0.4), items)

    cell("Strengths", swot.get("S", []), x1, y1)
    cell("Weaknesses", swot.get("W", []), x2, y1)
    cell("Opportunities", swot.get("O", []), x1, y2)
    cell("Threats", swot.get("T", []), x2, y2)
    return slide


def slide_ansoff(prs: Presentation, ansoff: Dict[str, List[str]]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _add_heading(slide, "Ansoff Matrix")
    grid_left, grid_top, grid_w, grid_h = MARGIN, Inches(1.5), W - 2*MARGIN, Inches(4.6)
    quads = _grid(slide, grid_left, grid_top, grid_w, grid_h)

    labels = [
        ("Market Penetration", ansoff.get("market_penetration", [])),
        ("Product Development", ansoff.get("product_development", [])),
        ("Market Development", ansoff.get("market_development", [])),
        ("Diversification", ansoff.get("diversification", [])),
    ]
    for (title, items), (l, t, w, h) in zip(labels, quads):
        title_box = slide.shapes.add_textbox(l + Inches(0.1), t + Inches(0.05), w - Inches(0.2), Inches(0.3))
        tf = title_box.text_frame; tf.clear(); p = tf.paragraphs[0]; r = p.add_run(); r.text = title; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = COLOR_PRIMARY
        _add_bullets(slide, l + Inches(0.1), t + Inches(0.45), w - Inches(0.2), h - Inches(0.6), items)

    _add_small_label(slide, "Existing Products → New Products", grid_left + grid_w/2 - Inches(1.2), grid_top - Inches(0.35),angle_deg=0)
    _add_small_label(slide, "New Markets → Existing Markets", grid_left - Inches(1.1), grid_top + grid_h/2 + Inches(0.05),angle_deg=270)
    return slide


def slide_benchmark(prs: Presentation, company: str, bench: Dict[str, Any]):
    table = bench.get("table") or []
    peers = bench.get("peers") or []
    if not table:
        return None
    cols = 2 + len(peers)  # Capability + company + peers
    rows = 1 + len(table)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _add_heading(slide, "Competitor Benchmark")

    left, top, width, height = MARGIN, Inches(1.2), W - 2*MARGIN, Inches(5.2)
    shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = shape.table

    # Header
    hdrs = ["Capability", company] + peers
    for j, h in enumerate(hdrs):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
        cell.fill.solid(); cell.fill.fore_color.rgb = COLOR_LIGHT

    # Body
    for i, row in enumerate(table, start=1):
        vals = [row.get("capability", "")] + [row.get(company, "")] + [row.get(p, "") for p in peers]
        for j, val in enumerate(vals):
            tbl.cell(i, j).text = str(val)

    return slide


def slide_recommendations(prs: Presentation, recs: List[Dict[str, Any]]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _add_heading(slide, "Top 5 Recommendations — Impact × Effort")

    grid_left, grid_top = MARGIN, Inches(2)
    grid_w, grid_h = W - 2*MARGIN, Inches(4.6)
    q = _grid(slide, grid_left, grid_top, grid_w, grid_h)

    # Quadrant labeling
    _add_small_label(slide, "Impact ->", grid_left - Inches(1.1), grid_top + grid_h - Inches(1.0), angle_deg=270)
    _add_small_label(slide, "Effort ->", grid_left + grid_w - Inches(0.8), grid_top + grid_h + Inches(0.05),angle_deg=0)

    # Place recs into quadrants
    #top = 0
    for idx, rec in enumerate((recs or [])[:10], start=1):
        title = rec.get("title", f"Rec {idx}")
        impact = int(rec.get("impact", 3))
        effort = int(rec.get("effort", 3))
        # Determine quadrant (simple thresholds)
        # TL (Q1): high impact (>=4), low effort (<=3)
        # TR (Q2): high impact, high effort (>3)
        # BL (Q3): low impact (<4), low effort (<=3)
        # BR (Q4): low impact, high effort (>3)
        if impact >= 4 and effort <= 3:
            quad = q[0]
        elif impact >= 4 and effort > 3:
            quad = q[1]
        elif impact < 4 and effort <= 3:
            quad = q[2]
        else:
            quad = q[3]
        l, t, w, h = quad
        left   = l + Inches(0.12)
        top    = t + Inches(0.12)
        width  = w - Inches(0.24)
        height = Inches(0.5)
        tol = Inches(0.05)
        gap = Inches(0.10)
        for shp in prs.slides[-1].shapes:
            #st.write(title, "left:", left, "shape left:", shp.left, "top:", top, "shape top:", shp.top)
            # horizontal overlap with our column
            if abs(left - shp.left) <= 1 and abs(top - shp.top) <= 1:
                top = shp.top + height
                #st.write(title, "left:", left, "shape left:", shp.left, "new top:", top, "shape top:", shp.top)
        box = prs.slides[-1].shapes.add_textbox(left, top, width, height)
        tf = box.text_frame; tf.clear(); p = tf.paragraphs[0]
        r = p.add_run(); r.text = f"{idx}. {title}"; r.font.size = BODY_SIZE; r.font.color.rgb = COLOR_DARK

    # Legend
    legend = slide.shapes.add_textbox(MARGIN, grid_top + grid_h + Inches(0.2), W - 2*MARGIN, Inches(0.6))
    tfl = legend.text_frame; tfl.clear()
    p = tfl.paragraphs[0]; r = p.add_run(); r.text = "Q1: Quick Wins   Q2: Strategic Bets   Q3: Fill-ins   Q4: Long Shots"; r.font.size = Pt(12); r.font.color.rgb = COLOR_MED
    return slide


def slide_appendix_json(prs: Presentation, title: str, text: str):
    # Break long text into multiple slides
    MAX_CHARS = 2000
    chunks = [text[i:i+MAX_CHARS] for i in range(0, len(text), MAX_CHARS)] or [text]
    for i, chunk in enumerate(chunks, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        _add_heading(slide, f"{title}{' (cont.)' if i>1 else ''}")
        box = slide.shapes.add_textbox(MARGIN, Inches(1.2), W - 2*MARGIN, Inches(5.5))
        tf = box.text_frame; tf.word_wrap = True; tf.clear()
        p = tf.paragraphs[0]; r = p.add_run(); r.text = chunk; r.font.name = "Courier New"; r.font.size = MONO_SIZE; r.font.color.rgb = COLOR_DARK

# ---------------------------- Orchestrator ----------------------------

def build_ppt_from_state(state: Dict[str, Any]) -> (BytesIO, str):
    """Return (pptx_bytes, filename) for download.
    Expects keys in `state`: company, product, frameworks, results, recs
    """
    company = (state.get("company") or "Company").strip()
    product = (state.get("product") or "Product").strip()
    results = state.get("results") or {}
    recs = state.get("recs") or []

    prs = Presentation()
    prs.slide_width, prs.slide_height = int(W), int(H)

    # Title
    date_str = datetime.now().strftime("%b %d, %Y")
    _add_title(prs, f"{product} × {company}", f"Strategy Snapshot — {date_str}")

    # Agenda
    slide_agenda(prs)

    # Executive Snapshot (basic heuristic based on SWOT + Ansoff presence)
    snapshot: List[str] = []
    if results.get("ind"):
        snapshot.append("Industry Analysis")
    swot = results.get("SWOT") or {}
    if any(swot.get(k) for k in ("S","W","O","T")):
        if swot.get("S"): snapshot.append(f"Strengths: {', '.join(swot['S'][:2])}")
        if swot.get("W"): snapshot.append(f"Weaknesses: {', '.join(swot['W'][:2])}")
        if swot.get("O"): snapshot.append(f"Opportunities: {', '.join(swot['O'][:2])}")
        if swot.get("T"): snapshot.append(f"Threats: {', '.join(swot['T'][:2])}")
    if results.get("Ansoff"):
        snapshot.append("Focus: Execute 1–2 high‑impact Ansoff plays next quarter.")
    if recs:
        snapshot.append(f"Top priority: {recs[0].get('title','First recommendation')}")
    slide_exec_snapshot(prs, snapshot[:6])

    # Industry Analysis
    Ind = results.get("ind") or {}
    if Ind:
        slide_ind(prs, Ind)
    
    # SWOT
    if swot:
        slide_swot(prs, swot)

    # Ansoff
    ansoff = results.get("Ansoff") or {}
    if ansoff:
        slide_ansoff(prs, ansoff)

    # Benchmark
    bench = results.get("Benchmark") or {}
    if bench.get("table"):
        slide_benchmark(prs, company, bench)

    # Recommendations
    if recs:
        slide_recommendations(prs, recs)

    # Appendix with raw JSON (trimmed)
    import json
    raw = json.dumps({
        "frameworks": state.get("frameworks", []),
        "results": results,
        "recs": recs,
    }, indent=2, ensure_ascii=False)
    slide_appendix_json(prs, "Appendix — Raw Analysis JSON", raw)

    # Serialize
    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)

    safe_company = company.replace(" ", "_")
    safe_product = product.replace(" ", "_")
    fname = f"{safe_company}_{safe_product}_{datetime.now().strftime('%Y%m%d')}_strategy.pptx"
    return bio, fname

# ---------------------------- Manual test ----------------------------
if __name__ == "__main__":  # pragma: no cover
    sample_state = {
        "company": "ACME Robotics",
        "product": "Industrial IoT Sensors",
        "frameworks": ["SWOT", "Ansoff", "Benchmark"],
        "results": {
            "SWOT": {"S": ["Reliable hardware", "OEM partners"], "W": ["Brand"], "O": ["Analytics"], "T": ["Price rivals"]},
            "Ansoff": {
                "market_penetration": ["Bundle add‑ons", "Loyalty pricing"],
                "market_development": ["Enter Canada"],
                "product_development": ["Managed calibration"],
                "diversification": ["Edge‑AI module"],
            },
            "Benchmark": {
                "peers": ["Rival A", "Rival B"],
                "table": [
                    {"capability": "Sensor accuracy", "ACME Robotics": "High", "Rival A": "High", "Rival B": "Medium"},
                    {"capability": "Cloud analytics", "ACME Robotics": "Medium", "Rival A": "High", "Rival B": "Low"},
                ],
            },
        },
        "recs": [
            {"title": "OEM Bundle Program", "impact": 5, "effort": 3},
            {"title": "Managed Calibration Add‑On", "impact": 4, "effort": 3},
            {"title": "Security Proof Pack", "impact": 3, "effort": 2},
        ],
    }
    bio, name = build_ppt_from_state(sample_state)
    with open(name, "wb") as f:
        f.write(bio.getvalue())
    print("Wrote", name)
